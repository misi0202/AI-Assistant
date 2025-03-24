from pptx import Presentation
import os, sys
import shutil
import numpy as np
from zhipuai import ZhipuAI
from sentence_transformers import SentenceTransformer, util
current_file_path = os.path.abspath(__file__)
parent_dir = os.path.dirname(current_file_path)
grandpa_dir = os.path.dirname(parent_dir)
sys.path.append(grandpa_dir)


from vectorstore.milvus import Milvus
from extract.extract_MinerU import add_Emedding2Chunk
def handlePPTFile(Course_name, partition_name, file_path):
    pdf_name = os.path.basename(file_path).split(".")[0]
    slides_text = extract_text_from_pptx(file_path)

    model_1 = SentenceTransformer("E:\\vue_pro\\ai_assistant\\ai_assistant_func\\model\\Conan")
    model_2 = ZhipuAI(api_key = "864eeb3324cb0bd34584e397a70caacf.jllKABCzmHYJPxfV") 

    chunks = []
    split_title = []
    split_content = []
    for slide, text in slides_text.items(): 
        # 排除内容过短的页
        if len(text) < 30:
            continue
        chunks.append([slide, text])
        split_title.append(slide)
        split_content.append(text)

    vector_1 = add_Emedding2Chunk(chunks, model_1, "Conan")
    vector_2 = add_Emedding2Chunk(chunks, model_2, "embedding-3")
    vector_1 = [np.array(vec, dtype=np.float32) for vec in vector_1]
    vector_2 = [np.array(vec, dtype=np.float32) for vec in vector_2]

    Course_vector = model_1.encode(Course_name)

    vectorstore = Milvus()
    # 创建collection
    vectorstore.create_chunk_collection()
    vectorstore.create_course_collection()

    vectorstore.create_partition(partition_name)
    vectorstore.insert_data(partition_name, split_title, split_content, vector_1, vector_2, Course_name, pdf_name)
    vectorstore.insert_course_data(Course_name, partition_name, pdf_name, Course_vector)

    return True

def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text[f"Slide {i}"] = "\n".join(slide_text)

    return slides_text

if __name__ == "__main__":
    # 读取 PPTX 文件并按页显示文字
    pptx_file = "E:\\web_download\\book\\Week 4 朴素贝叶斯.pptx"  # 替换为你的 PPTX 文件路径
    slides_text = extract_text_from_pptx(pptx_file)

    # 打印每页的文字内容
    for slide, text in slides_text.items():
        print(f"{slide}:\n{text}\n{'-'*40}")

