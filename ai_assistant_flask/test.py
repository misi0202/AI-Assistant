from pptx import Presentation
from pathlib import Path
def extract_text_from_pptx(pptx_path):
    prs = Presentation(pptx_path)
    slides_text = {}

    for i, slide in enumerate(prs.slides, start=1):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text[f"Slide {i}"] = "\n".join(slide_text)

    return slides_text

# 读取 PPTX 文件并按页显示文字
pptx_file = "E:\\web_download\\book\\Week 4 朴素贝叶斯.pptx"  # 替换为你的 PPTX 文件路径
slides_text = extract_text_from_pptx(pptx_file)

# 打印每页的文字内容
for slide, text in slides_text.items():
    print(f"{slide}:\n{text}\n{'-'*40}")
print(Path(pptx_file).suffix.lower())
